"""Gera apresentacao.pptx (16:9) a partir do conteudo do ROTEIRO.md.

Uso:
    python gerar_slides.py

Requer: python-pptx (pip install python-pptx).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# Identidade visual (mesma dos relatorios)
FONT = "Calibri"
TITLE_FONT = "Calibri"
MONO_FONT = "Consolas"

COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)
COLOR_TEXT = RGBColor(0x22, 0x22, 0x22)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_ACCENT = RGBColor(0x2E, 0x75, 0xB6)
COLOR_CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)
COLOR_OK = RGBColor(0x2E, 0x8B, 0x57)

# 16:9 = 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_X = Inches(0.6)
MARGIN_Y = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_X * 2


# --------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------- #

def add_textbox(slide, left, top, width, height, text, *, font_size=18,
                bold=False, italic=False, color=None, font_name=None,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name or FONT
    if color is not None:
        run.font.color.rgb = color
    return tb


def add_multiline(slide, left, top, width, height, lines, *, font_size=18,
                  bold=False, color=None, font_name=None,
                  bullet=False, line_spacing=1.15):
    """lines pode ser list[str] ou list[dict] com chaves: text, size, bold,
    italic, color, font, bullet, indent."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        if isinstance(line, str):
            line = {"text": line}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        p.space_after = Pt(line.get("space_after", 4))
        is_bullet = line.get("bullet", bullet)
        if is_bullet:
            _set_bullet(p, level=line.get("indent", 0))
        run = p.add_run()
        run.text = line["text"]
        run.font.size = Pt(line.get("size", font_size))
        run.font.bold = line.get("bold", bold)
        run.font.italic = line.get("italic", False)
        run.font.name = line.get("font", font_name or FONT)
        clr = line.get("color", color)
        if clr is not None:
            run.font.color.rgb = clr
    return tb


def _set_bullet(paragraph, level=0):
    """Adiciona marcador • ao paragrafo."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("indent", str(-228600))  # -0.25in
    pPr.set("marL", str(228600 + level * 228600))
    pPr.set("lvl", str(level))
    # remove qualquer buChar/buNone existente
    for tag in ("a:buChar", "a:buNone", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "•")


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if line_color is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line_color
    rect.shadow.inherit = False
    # remove texto default
    rect.text_frame.text = ""
    return rect


def add_section_header(slide, eyebrow, title):
    """Cabecalho padrao: faixa colorida na esquerda + eyebrow + titulo grande."""
    # Faixa lateral
    add_filled_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, COLOR_PRIMARY)
    # Eyebrow (texto pequeno por cima)
    if eyebrow:
        add_textbox(
            slide, MARGIN_X, Inches(0.45), CONTENT_W, Inches(0.4),
            eyebrow, font_size=14, bold=True, color=COLOR_ACCENT,
        )
    add_textbox(
        slide, MARGIN_X, Inches(0.8), CONTENT_W, Inches(0.9),
        title, font_size=32, bold=True, color=COLOR_PRIMARY,
    )


def add_code_panel(slide, left, top, width, height, code, *, font_size=14):
    """Painel de codigo monoespacado com fundo cinza claro."""
    add_filled_rect(slide, left, top, width, height, COLOR_CODE_BG)
    pad = Inches(0.18)
    lines = code.split("\n")
    line_objs = []
    for ln in lines:
        line_objs.append({
            "text": ln if ln else " ",
            "font": MONO_FONT,
            "size": font_size,
            "color": COLOR_TEXT,
            "space_after": 0,
        })
    add_multiline(
        slide, left + pad, top + pad, width - 2 * pad, height - 2 * pad,
        line_objs, line_spacing=1.05,
    )


def add_footer(slide, page_num, total, speaker=None):
    """Rodape com numero do slide e quem fala."""
    pad = Inches(0.25)
    right = SLIDE_W - pad - Inches(0.6)
    add_textbox(
        slide, right, SLIDE_H - Inches(0.4), Inches(0.6), Inches(0.3),
        f"{page_num}/{total}", font_size=10, color=COLOR_MUTED,
        align=PP_ALIGN.RIGHT,
    )
    if speaker:
        add_textbox(
            slide, MARGIN_X, SLIDE_H - Inches(0.4), Inches(6.0), Inches(0.3),
            f"Apresenta: {speaker}", font_size=10, color=COLOR_MUTED,
        )


# --------------------------------------------------------------------- #
# Construcao da apresentacao
# --------------------------------------------------------------------- #

TOTAL_SLIDES = 15


def new_blank_slide(prs):
    blank = prs.slide_layouts[6]  # layout em branco
    return prs.slides.add_slide(blank)


def slide_capa(prs):
    s = new_blank_slide(prs)
    # Fundo solido + faixa colorida no topo
    add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.8), COLOR_PRIMARY)
    # Eyebrow institucional
    add_textbox(
        s, MARGIN_X, Inches(0.5), CONTENT_W, Inches(0.4),
        "Universidade Federal da Paraíba — Centro de Informática",
        font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF),
    )
    add_textbox(
        s, MARGIN_X, Inches(0.9), CONTENT_W, Inches(0.5),
        "Construção de Compiladores 1 — Prof. Andrei de Araújo Formiga",
        font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF), italic=True,
    )
    # Titulo grande
    add_textbox(
        s, MARGIN_X, Inches(2.6), CONTENT_W, Inches(1.2),
        "Marco 1 — Compilador EC1", font_size=44, bold=True,
        color=COLOR_PRIMARY, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, MARGIN_X, Inches(3.8), CONTENT_W, Inches(0.7),
        "Análise Léxica · Análise Sintática · Geração de Código",
        font_size=20, italic=True, color=COLOR_MUTED,
        align=PP_ALIGN.CENTER,
    )
    # Integrantes
    add_textbox(
        s, MARGIN_X, Inches(5.2), CONTENT_W, Inches(0.4),
        "Integrantes do grupo", font_size=14, bold=True,
        color=COLOR_ACCENT, align=PP_ALIGN.CENTER,
    )
    add_multiline(
        s, MARGIN_X, Inches(5.7), CONTENT_W, Inches(1.4),
        [
            {"text": "Davi Alves Rodrigues", "align": PP_ALIGN.CENTER, "size": 16},
            {"text": "Gabriel Barbosa Ribeiro de Oliveira",
             "align": PP_ALIGN.CENTER, "size": 16},
            {"text": "João Vitor Sampaio Costa",
             "align": PP_ALIGN.CENTER, "size": 16},
            {"text": "Nathan Meira Nóbrega",
             "align": PP_ALIGN.CENTER, "size": 16},
        ],
        line_spacing=1.2,
    )
    add_footer(s, 1, TOTAL_SLIDES, "Davi")


def slide_linguagem_ec1(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Contexto", "A linguagem EC1")
    # corpo: gramatica + exemplos lado a lado
    col_w = (CONTENT_W - Inches(0.6)) // 2
    add_textbox(
        s, MARGIN_X, Inches(2.0), col_w, Inches(0.4),
        "Gramática", font_size=16, bold=True, color=COLOR_ACCENT,
    )
    add_code_panel(
        s, MARGIN_X, Inches(2.5), col_w, Inches(2.6),
        "<programa>  ::= <expressao>\n"
        "<expressao> ::= <literal>\n"
        "             | '(' <expressao> <operador> <expressao> ')'\n"
        "<operador>  ::= '+' | '-' | '*' | '/'\n"
        "<literal>   ::= <digito>+",
        font_size=14,
    )
    right_x = MARGIN_X + col_w + Inches(0.6)
    add_textbox(
        s, right_x, Inches(2.0), col_w, Inches(0.4),
        "Exemplos válidos", font_size=16, bold=True, color=COLOR_ACCENT,
    )
    add_code_panel(
        s, right_x, Inches(2.5), col_w, Inches(2.6),
        "333\n"
        "(6 * 7)\n"
        "(33 + (912 * 11))\n"
        "((427 / 7) + (11 * (231 + 5)))",
        font_size=15,
    )
    add_multiline(
        s, MARGIN_X, Inches(5.5), CONTENT_W, Inches(1.5),
        [
            {"text": "Toda operação é parentizada → sem regras de precedência.",
             "bullet": True},
            {"text": "Operandos são sempre constantes inteiras.",
             "bullet": True},
            {"text": "A linguagem é mínima, mas exibe o pipeline completo de um compilador.",
             "bullet": True},
        ],
        font_size=16,
    )
    add_footer(s, 2, TOTAL_SLIDES, "Davi")


def slide_pipeline(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Visão geral", "Pipeline do compilador EC1")

    # Diagrama: 5 caixas horizontais (3 etapas + entradas/saidas)
    y = Inches(3.0)
    h = Inches(1.2)
    gap = Inches(0.25)
    n = 5
    total_w = CONTENT_W
    box_w = (total_w - gap * (n - 1)) // n
    x = MARGIN_X
    items = [
        ("Código\nfonte (.ec1)", COLOR_MUTED, RGBColor(0xFF, 0xFF, 0xFF)),
        ("Lexer\n(Atividade 04)", COLOR_PRIMARY, RGBColor(0xFF, 0xFF, 0xFF)),
        ("Parser → AST\n(Atividade 05)", COLOR_PRIMARY, RGBColor(0xFF, 0xFF, 0xFF)),
        ("Codegen\n(Atividade 06)", COLOR_PRIMARY, RGBColor(0xFF, 0xFF, 0xFF)),
        ("Assembly\n(.s) → exec.", COLOR_MUTED, RGBColor(0xFF, 0xFF, 0xFF)),
    ]
    for i, (label, bg, fg) in enumerate(items):
        rect = add_filled_rect(s, x, y, box_w, h, bg)
        tf = rect.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.name = FONT
        run.font.color.rgb = fg
        if i < n - 1:
            # seta entre caixas
            arrow_x = x + box_w + Inches(0.02)
            arrow_w = gap - Inches(0.04)
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, y + h // 2 - Inches(0.15),
                arrow_w, Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.fill.background()
        x += box_w + gap

    # legenda embaixo
    add_multiline(
        s, MARGIN_X, Inches(5.0), CONTENT_W, Inches(1.6),
        [
            {"text": "Lexer: caractere → token (NUMERO, PAREN_ESQ, …).",
             "bullet": True},
            {"text": "Parser: tokens → árvore sintática (Const, OpBin).",
             "bullet": True},
            {"text": "Codegen: árvore → assembly x86-64 que coloca o resultado em %rax.",
             "bullet": True},
        ],
        font_size=16,
    )
    add_footer(s, 3, TOTAL_SLIDES, "Davi")


def slide_lex_intro(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 04", "Análise Léxica")
    add_multiline(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(0.8),
        [{"text": "Varre a string fonte e agrupa caracteres em tokens.",
          "size": 18}],
    )
    add_textbox(
        s, MARGIN_X, Inches(2.9), CONTENT_W, Inches(0.4),
        "Tipos de token (TipoToken)", font_size=16, bold=True,
        color=COLOR_ACCENT,
    )
    add_code_panel(
        s, MARGIN_X, Inches(3.4), Inches(6.0), Inches(3.5),
        "NUMERO      sequência de dígitos\n"
        "PAREN_ESQ   (\n"
        "PAREN_DIR   )\n"
        "SOMA        +\n"
        "SUB         -\n"
        "MULT        *\n"
        "DIV         /\n"
        "EOF         fim da entrada",
        font_size=14,
    )
    add_multiline(
        s, Inches(7.2), Inches(3.4), Inches(5.6), Inches(3.5),
        [
            {"text": "Espaços, tabs, \\n e \\r são descartados.",
             "bullet": True},
            {"text": "Números são agrupados num único token com o lexema completo.",
             "bullet": True},
            {"text": "Caracteres simples → lookup direto em dict (O(1)).",
             "bullet": True},
            {"text": "Extensão opcional: comentários de linha iniciados por #.",
             "bullet": True},
        ],
        font_size=15,
    )
    add_footer(s, 4, TOTAL_SLIDES, "Nathan")


def slide_lex_token(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 04", "Estrutura do Token e API do lexer")
    add_multiline(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(0.6),
        [{"text": "Cada token guarda tipo, lexema e posição no arquivo "
                  "— essencial para mensagens de erro.", "size": 16}],
    )
    add_code_panel(
        s, MARGIN_X, Inches(2.8), Inches(6.0), Inches(2.3),
        "@dataclass(frozen=True)\n"
        "class Token:\n"
        "    tipo: TipoToken\n"
        "    lexema: str\n"
        "    posicao: int",
        font_size=14,
    )
    add_textbox(
        s, Inches(7.2), Inches(2.8), Inches(5.6), Inches(0.4),
        "API pública de AnalisadorLexico", font_size=16, bold=True,
        color=COLOR_ACCENT,
    )
    add_code_panel(
        s, Inches(7.2), Inches(3.2), Inches(5.6), Inches(1.9),
        "proximo_token() -> Token\n"
        "    consome e devolve\n"
        "\n"
        "olhar_proximo() -> Token\n"
        "    peek (não consome)",
        font_size=14,
    )
    add_multiline(
        s, MARGIN_X, Inches(5.5), CONTENT_W, Inches(1.5),
        [
            {"text": "Cursor interno avança caractere a caractere.",
             "bullet": True},
            {"text": "Buffer de 1 token suporta o peek sem perder a posição.",
             "bullet": True},
            {"text": "Erro léxico vira exceção com posição e ASCII do caractere ofensor.",
             "bullet": True},
        ],
        font_size=15,
    )
    add_footer(s, 5, TOTAL_SLIDES, "Nathan")


def slide_lex_exemplo(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 04", "Exemplo: (33 + (912 * 11))")
    add_textbox(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(0.4),
        "Entrada", font_size=14, bold=True, color=COLOR_ACCENT,
    )
    add_code_panel(
        s, MARGIN_X, Inches(2.4), CONTENT_W, Inches(0.7),
        "(33 + (912 * 11))",
        font_size=18,
    )
    add_textbox(
        s, MARGIN_X, Inches(3.4), CONTENT_W, Inches(0.4),
        "Tokens produzidos", font_size=14, bold=True, color=COLOR_ACCENT,
    )
    add_code_panel(
        s, MARGIN_X, Inches(3.8), CONTENT_W, Inches(2.1),
        '<ParenEsq, "(", 0>    <Numero, "33", 1>    <Soma, "+", 4>\n'
        '<ParenEsq, "(", 6>    <Numero, "912", 7>   <Mult, "*", 11>\n'
        '<Numero, "11", 13>    <ParenDir, ")", 15>  <ParenDir, ")", 16>',
        font_size=14,
    )
    add_multiline(
        s, MARGIN_X, Inches(6.1), CONTENT_W, Inches(1.0),
        [
            {"text": "Posição é a do caractere de início no arquivo original "
                     "(não o índice na sequência).", "bullet": True},
            {"text": "43 testes em unittest cobrem tipos, lexemas, posições "
                     "e erros léxicos.", "bullet": True},
        ],
        font_size=14,
    )
    add_footer(s, 6, TOTAL_SLIDES, "Nathan")


def slide_parse_intro(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 05", "Análise Sintática")
    add_multiline(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(0.8),
        [{"text": "Transforma a sequência de tokens em uma árvore sintática "
                  "que captura a estrutura da expressão.", "size": 18}],
    )
    add_textbox(
        s, MARGIN_X, Inches(2.9), CONTENT_W, Inches(0.4),
        "Algoritmo: descendente recursivo",
        font_size=16, bold=True, color=COLOR_ACCENT,
    )
    add_multiline(
        s, MARGIN_X, Inches(3.4), CONTENT_W, Inches(3.0),
        [
            {"text": "Uma função por não-terminal da gramática.",
             "bullet": True},
            {"text": "Decisão por lookahead de 1 token (consome e ramifica).",
             "bullet": True},
            {"text": "Recursão natural: a regra de <expressao> referencia a si própria.",
             "bullet": True},
            {"text": "Erro sintático vira exceção com posição e mensagem "
                     "(esperado X, encontrado Y).", "bullet": True},
        ],
        font_size=16,
    )
    add_footer(s, 7, TOTAL_SLIDES, "João Vitor")


def slide_parse_codigo(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 05", "_analisa_exp() — pseudocódigo")
    add_code_panel(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(3.6),
        "def _analisa_exp(self) -> Exp:\n"
        "    tok = self._lex.proximo_token()\n"
        "    if tok.tipo == NUMERO:\n"
        "        return Const(int(tok.lexema))\n"
        "    if tok.tipo == PAREN_ESQ:\n"
        "        esq = self._analisa_exp()\n"
        "        op  = self._analisa_operador()\n"
        "        dir = self._analisa_exp()\n"
        "        self._verifica_token(PAREN_DIR)\n"
        "        return OpBin(op, esq, dir)\n"
        "    raise ErroSintatico(...)",
        font_size=15,
    )
    add_multiline(
        s, MARGIN_X, Inches(5.8), CONTENT_W, Inches(1.2),
        [{"text": "analisa_programa() chama _analisa_exp() e ainda exige EOF — "
                  "captura entradas como `(6 * 7) 42` que têm lixo após a raiz.",
          "bullet": True, "size": 15}],
    )
    add_footer(s, 8, TOTAL_SLIDES, "João Vitor")


def slide_ast(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 05", "Árvore de sintaxe abstrata")
    add_code_panel(
        s, MARGIN_X, Inches(2.0), Inches(6.4), Inches(3.6),
        "class Exp: ...  # abstrata\n"
        "\n"
        "@dataclass(frozen=True)\n"
        "class Const(Exp):\n"
        "    valor: int\n"
        "\n"
        "@dataclass(frozen=True)\n"
        "class OpBin(Exp):\n"
        "    op: Op\n"
        "    esq: Exp\n"
        "    dir: Exp",
        font_size=15,
    )
    add_textbox(
        s, Inches(7.2), Inches(2.0), Inches(5.6), Inches(0.4),
        "AST de (33 + (912 * 11))", font_size=14, bold=True,
        color=COLOR_ACCENT,
    )
    add_code_panel(
        s, Inches(7.2), Inches(2.4), Inches(5.6), Inches(3.2),
        "      +\n"
        "     / \\\n"
        "   33   *\n"
        "       / \\\n"
        "     912  11",
        font_size=18,
    )
    add_multiline(
        s, MARGIN_X, Inches(5.9), CONTENT_W, Inches(1.2),
        [{"text": "frozen=True dá __eq__ por valor de graça — comparações "
                  "diretas nos testes (self.assertEqual(arvore, esperado)).",
          "bullet": True, "size": 14}],
    )
    add_footer(s, 9, TOTAL_SLIDES, "João Vitor")


def slide_interpretador(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 05", "Interpretador (tree-walking)")
    add_multiline(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(0.8),
        [{"text": "Para validar que a AST captura a expressão certa, "
                  "implementamos um interpretador como método na própria árvore.",
          "size": 16}],
    )
    add_code_panel(
        s, MARGIN_X, Inches(2.9), Inches(6.4), Inches(3.0),
        "class OpBin(Exp):\n"
        "    def avaliar(self) -> int:\n"
        "        ve = self.esq.avaliar()\n"
        "        vd = self.dir.avaliar()\n"
        "        if self.op is Op.SOMA: return ve + vd\n"
        "        if self.op is Op.SUB:  return ve - vd\n"
        "        if self.op is Op.MULT: return ve * vd\n"
        "        if self.op is Op.DIV:  return int(ve / vd)",
        font_size=14,
    )
    add_textbox(
        s, Inches(7.2), Inches(2.9), Inches(5.6), Inches(0.4),
        "Demo no terminal", font_size=14, bold=True, color=COLOR_ACCENT,
    )
    add_code_panel(
        s, Inches(7.2), Inches(3.3), Inches(5.6), Inches(2.0),
        "$ python ec1.py exemplos/valido3.ec1\n10065",
        font_size=15,
    )
    add_multiline(
        s, MARGIN_X, Inches(6.1), CONTENT_W, Inches(1.0),
        [{"text": "Esse interpretador vai servir como gabarito para validar "
                  "o gerador de código da Atividade 06.",
          "bullet": True, "size": 14, "color": COLOR_OK}],
    )
    add_footer(s, 10, TOTAL_SLIDES, "João Vitor")


def slide_codegen_problema(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 06", "O problema da geração de código")
    add_multiline(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(1.6),
        [
            {"text": "Como traduzir (A op B) para assembly sem perder o resultado "
                     "de A enquanto se calcula B?", "size": 17},
            {"text": "Reservar um registrador extra (RBX) funciona para "
                     "(7 + 11), mas quebra para (7 + (3 + 8)) — RBX é "
                     "sobrescrito pelo aninhamento.",
             "size": 15, "color": COLOR_MUTED, "italic": True},
        ],
    )
    add_code_panel(
        s, MARGIN_X, Inches(4.0), CONTENT_W, Inches(2.6),
        "# tentativa ingenua com RBX para (7 + (3 + 8))\n"
        "mov  $7, %rax\n"
        "mov  %rax, %rbx     # guarda 7\n"
        "mov  $3, %rax\n"
        "mov  %rax, %rbx     # ← ALTERA 7 PARA 3!\n"
        "mov  $8, %rax\n"
        "add  %rbx, %rax\n"
        "add  %rbx, %rax     # 7 perdido — resultado errado",
        font_size=14,
    )
    add_footer(s, 11, TOTAL_SLIDES, "Gabriel")


def slide_codegen_esquema(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 06", "Esquema da pilha (opção 2)")
    add_multiline(
        s, MARGIN_X, Inches(2.0), Inches(7.0), Inches(4.5),
        [
            {"text": "Para OpBin(op, esq, dir):", "size": 15, "bold": True,
             "color": COLOR_ACCENT},
            {"text": "1. gera código de dir   → %rax", "size": 15},
            {"text": "2. push %rax", "size": 15},
            {"text": "3. gera código de esq   → %rax", "size": 15},
            {"text": "4. pop %rbx", "size": 15},
            {"text": "5. <instr> %rbx, %rax    → %rax = esq op dir", "size": 15},
            {"text": "", "size": 8},
            {"text": "A ordem invertida já deixa esq em RAX e dir em RBX — "
                     "natural para sub e idiv, sem truques.",
             "size": 14, "italic": True, "color": COLOR_OK},
        ],
        line_spacing=1.25,
    )
    add_textbox(
        s, Inches(9.4), Inches(2.0), Inches(3.4), Inches(0.4),
        "Mapeamento dos operadores", font_size=14, bold=True,
        color=COLOR_ACCENT,
    )
    add_code_panel(
        s, Inches(9.4), Inches(2.4), Inches(3.4), Inches(4.0),
        "SOMA  →  add  %rbx, %rax\n"
        "SUB   →  sub  %rbx, %rax\n"
        "MULT  →  imul %rbx, %rax\n"
        "DIV   →  cqo\n"
        "         idiv %rbx\n"
        "\n"
        "(cqo estende o sinal\n"
        "de rax em rdx, exigido\n"
        "antes de idiv)",
        font_size=13,
    )
    add_footer(s, 12, TOTAL_SLIDES, "Gabriel")


def slide_demo(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 06", "Demo: (33 + (912 * 11)) → .s")
    add_textbox(
        s, MARGIN_X, Inches(2.0), Inches(6.2), Inches(0.4),
        "$ python compec1.py exemplos/valido3.ec1",
        font_size=14, font_name=MONO_FONT, color=COLOR_ACCENT,
    )
    add_code_panel(
        s, MARGIN_X, Inches(2.4), Inches(6.2), Inches(4.6),
        "    .section .text\n"
        "    .globl _start\n"
        "_start:\n"
        "    mov $11, %rax     # dir.dir\n"
        "    push %rax\n"
        "    mov $912, %rax    # dir.esq\n"
        "    pop %rbx\n"
        "    imul %rbx, %rax   # 912*11 = 10032\n"
        "    push %rax\n"
        "    mov $33, %rax     # esq\n"
        "    pop %rbx\n"
        "    add %rbx, %rax    # 33+10032 = 10065\n"
        "    call imprime_num\n"
        "    call sair\n"
        '    .include \"runtime.s\"',
        font_size=12,
    )
    add_textbox(
        s, Inches(7.5), Inches(2.0), Inches(5.3), Inches(0.4),
        "Trace do esquema", font_size=14, bold=True, color=COLOR_ACCENT,
    )
    add_multiline(
        s, Inches(7.5), Inches(2.4), Inches(5.3), Inches(4.6),
        [
            {"text": "1. (912 * 11) é resolvida primeiro — imul deixa 10032 em %rax.",
             "bullet": True, "size": 13},
            {"text": "2. 10032 vai pra pilha.", "bullet": True, "size": 13},
            {"text": "3. 33 vai pra %rax.", "bullet": True, "size": 13},
            {"text": "4. pop traz 10032 para %rbx.", "bullet": True, "size": 13},
            {"text": "5. add → 10065 em %rax.", "bullet": True, "size": 13},
            {"text": "6. imprime_num imprime e sair termina o programa.",
             "bullet": True, "size": 13},
            {"text": "", "size": 8},
            {"text": "Pilha sempre balanceada: cada push tem seu pop.",
             "size": 13, "color": COLOR_OK, "italic": True},
        ],
        line_spacing=1.25,
    )
    add_footer(s, 13, TOTAL_SLIDES, "Gabriel")


def slide_testes(prs):
    s = new_blank_slide(prs)
    add_section_header(s, "Atividade 06", "Validação por equivalência semântica")
    add_multiline(
        s, MARGIN_X, Inches(2.0), CONTENT_W, Inches(0.8),
        [{"text": "15 testes em 6 classes. O destaque é o teste de equivalência:",
          "size": 16}],
    )
    add_code_panel(
        s, MARGIN_X, Inches(2.9), CONTENT_W, Inches(2.3),
        "# Para cada programa EC1 de teste:\n"
        "arvore   = analisar(fonte)\n"
        "esperado = arvore.avaliar()                 # interpretador (Ativ. 05)\n"
        "obtido   = simular(gerar_codigo(arvore))    # simulador da pilha\n"
        "assert obtido == esperado",
        font_size=14,
    )
    add_multiline(
        s, MARGIN_X, Inches(5.4), CONTENT_W, Inches(1.6),
        [
            {"text": "O simulador (~30 linhas) entende exatamente as instruções "
                     "que o codegen emite — mov, push, pop, add, sub, imul, cqo, idiv.",
             "bullet": True, "size": 14},
            {"text": "Se simulador e interpretador batem para 16 programas, "
                     "o gerador é semanticamente correto sem precisar montar/linkar.",
             "bullet": True, "size": 14, "color": COLOR_OK},
        ],
    )
    add_footer(s, 14, TOTAL_SLIDES, "Gabriel")


def slide_encerramento(prs):
    s = new_blank_slide(prs)
    add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.5), COLOR_PRIMARY)
    add_textbox(
        s, MARGIN_X, Inches(0.5), CONTENT_W, Inches(0.8),
        "Marco 1 — Encerramento", font_size=32, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
    )
    add_multiline(
        s, MARGIN_X, Inches(2.2), CONTENT_W, Inches(2.5),
        [
            {"text": "Compilador EC1 completo:", "size": 22, "bold": True,
             "color": COLOR_PRIMARY, "align": PP_ALIGN.CENTER},
            {"text": "lexer · parser · interpretador · gerador de código x86-64",
             "size": 18, "italic": True, "color": COLOR_MUTED,
             "align": PP_ALIGN.CENTER},
        ],
        line_spacing=1.3,
    )
    add_textbox(
        s, MARGIN_X, Inches(4.5), CONTENT_W, Inches(0.6),
        "Repositório no GitHub", font_size=14, bold=True,
        color=COLOR_ACCENT, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, MARGIN_X, Inches(5.0), CONTENT_W, Inches(0.6),
        "github.com/gabrielbribeiroo/CompilerConstruction_UFPB",
        font_size=18, font_name=MONO_FONT, color=COLOR_TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, MARGIN_X, Inches(6.3), CONTENT_W, Inches(0.5),
        "Obrigado!", font_size=26, bold=True,
        color=COLOR_PRIMARY, align=PP_ALIGN.CENTER,
    )
    add_footer(s, 15, TOTAL_SLIDES)


# --------------------------------------------------------------------- #

def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)                # 1   Davi
    slide_linguagem_ec1(prs)       # 2   Davi
    slide_pipeline(prs)            # 3   Davi
    slide_lex_intro(prs)           # 4   Gabriel
    slide_lex_token(prs)           # 5   Gabriel
    slide_lex_exemplo(prs)         # 6   Gabriel
    slide_parse_intro(prs)         # 7   Joao
    slide_parse_codigo(prs)        # 8   Joao
    slide_ast(prs)                 # 9   Joao
    slide_interpretador(prs)       # 10  Joao
    slide_codegen_problema(prs)    # 11  Nathan
    slide_codegen_esquema(prs)     # 12  Nathan
    slide_demo(prs)                # 13  Nathan
    slide_testes(prs)              # 14  Nathan
    slide_encerramento(prs)        # 15  encerramento

    prs.save(path)


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(__file__), "apresentacao.pptx")
    build(out)
    print(f"gerado: {out}")
